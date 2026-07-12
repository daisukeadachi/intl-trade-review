#!/usr/bin/env python3
"""Extract slide text (incl. tables and speaker notes) from .pptx lectures.

Used to draft review prompts from the course slides.

Usage:  python3 scripts/extract_pptx.py <slides_dir> [out_dir]
Deps:   pip install python-pptx
"""

import sys
from pathlib import Path

from pptx import Presentation


def shape_text(shape):
    parts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs)
            if t.strip():
                parts.append("  " * para.level + t.strip())
    if shape.shape_type == 6:  # group
        for s in shape.shapes:
            parts.extend(shape_text(s))
    if shape.has_table:
        for row in shape.table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append(" | ".join(cells))
    return parts


def main() -> int:
    if len(sys.argv) < 2:
        print(__doc__)
        return 1
    slides_dir = Path(sys.argv[1])
    out_dir = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("pptx_text")
    out_dir.mkdir(exist_ok=True)

    for f in sorted(slides_dir.glob("*.pptx")):
        prs = Presentation(f)
        lines = [f"##### {f.name} #####"]
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                try:
                    lines.extend(shape_text(shape))
                except Exception as e:  # noqa: BLE001 - keep going per shape
                    lines.append(f"[shape error: {e}]")
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"[NOTES] {notes[:500]}")
        out = out_dir / (f.stem + ".txt")
        out.write_text("\n".join(lines), encoding="utf-8")
        print(f"{f.name} -> {out}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
